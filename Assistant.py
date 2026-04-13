import os
import json
import re
import subprocess
import tempfile
import shutil
import uuid
import time
import webbrowser
import requests
import base64
import cv2
from pathlib import Path
from rapidfuzz import fuzz
from ddgs import DDGS
from datetime import datetime
from urllib.parse import quote
from requests.exceptions import ReadTimeout
from PIL import Image
import io
import ast 
import re
import threading
import queue

# File format support
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except ImportError:
    pdfminer_extract_text = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import fitz
except ImportError:
    fitz = None

try:
    import pytesseract
    def _find_tesseract_cmd():
        candidates = [
            shutil.which("tesseract"),
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        ]
        for candidate in candidates:
            if candidate and os.path.isfile(candidate):
                return candidate
        return None

    TESSERACT_CMD = _find_tesseract_cmd()
    POPPLER_PATHS = [
        r"C:\Program Files\poppler-25.12.0\Library\bin",
        r"C:\Program Files\poppler-22.12.0\Library\bin",
        r"C:\Users\Lenovo\Downloads\poppler-25.12.0\poppler-25.12.0\Library\bin",
        r"C:\Users\Lenovo\Downloads\Release-25.12.0-0\Library\bin",
    ]

    def _find_poppler_path():
        for candidate in POPPLER_PATHS:
            if candidate and os.path.isdir(candidate):
                return candidate
        return None

    POPPLER_CMD = _find_poppler_path()

    if TESSERACT_CMD:
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
    else:
        TESSERACT_CMD = None
except ImportError:
    pytesseract = None
    TESSERACT_CMD = None
    POPPLER_CMD = None

from app_automation import handle_app_command
from voice_response import speak_response, speak_immediate, toggle_voice, is_voice_enabled
from enhanced_intent import detect_enhanced_intent, separate_question_action
from tone_manager import get_tone_manager, set_ai_tone, get_current_tone_name, format_with_tone, get_tone_system_prompt

# =========================
# CONFIG
# =========================
ASSISTANT_NAME = "Jarvis"
MEMORY_FILE = "memory.json"
AI_MODEL = "qwen2.5:7b"  # Fast, efficient, less internet-dependent
VISION_MODEL = "llava:latest"  # Vision model for speed/quality balance
# CAD sub-assistant (OpenSCAD JSON). Use a capable model if available, e.g. "qwen2.5:7b" or "llama3.1:8b"
CAD_MODEL = "llama3:latest"  # Fast CAD model with good local knowledge
OLLAMA_URL = "http://localhost:11434/api/chat"
AI_TIMEOUT = 10  # Faster timeout for quicker responses
AI_RETRY_TIMEOUT = 30
AI_DOC_MODEL = "qwen2.5:7b"  # Fast, capable model for document review and analysis
AI_DOC_TIMEOUT = 45  # Longer timeout for document review and analysis
AI_DOC_RETRY_TIMEOUT = 90
CAD_TIMEOUT = 45
CAD_RETRY_TIMEOUT = 90
ENABLE_CAMERA = True  # Set to False to disable camera

# STL previews (compiled from OpenSCAD via CLI; served by web_app)
STL_PREVIEW_CACHE = {}
STL_CACHE_MAX_ITEMS = 40
STL_CACHE_TTL_SEC = 3600
OPENSCAD_COMPILE_TIMEOUT = 120
CAMERA_INDEX = 0  # Default camera (0 = primary webcam)
ENABLE_VOICE_RESPONSE = True  # Set to False to disable voice responses
ENABLE_TONE_SYSTEM = True  # Set to False to disable tone customization
ENABLE_CONTINUOUS_IMPROVEMENT = True  # Set to False to disable autonomous self-improvement
IMPROVEMENT_INTERVAL = 3600  # Seconds between self-improvement checks (1 hour)
IMPROVEMENT_LOG_FILE = "improvements.json"  # Log of applied improvements

# =========================
# MEMORY
# =========================
def load_memory():
    if os.path.exists(MEMORY_FILE):
        try:
            with open(MEMORY_FILE, "r") as f:
                return json.load(f)
        except:
            return {}
    return {}

memory = load_memory()
http = requests.Session()
_ROBLOX_CACHE = None

# =========================
# FILE EXTRACTION UTILITIES
# =========================
def _extract_text_from_pdf_bytes(file_path, min_length=20, max_chunks=50):
    try:
        with open(file_path, 'rb') as f:
            raw_bytes = f.read()
        segments = re.findall(br'[\t\n\r\x20-\x7E]{%d,}' % min_length, raw_bytes)
        if segments:
            extracted = "\n".join(
                segment.decode('latin1', errors='replace').strip()
                for segment in segments[:max_chunks]
                if segment.strip()
            )
            return extracted.strip()
    except Exception:
        pass
    return ""


def extract_pdf_text(file_path, max_chars=50000):
    """Extract text from PDF files with multiple fallback methods"""
    import time
    
    # First, validate the file
    if not os.path.exists(file_path):
        return None, f"File does not exist: {file_path}"
    
    try:
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            return None, "PDF file is empty (0 bytes)"
    except Exception as e:
        return None, f"Cannot read file: {str(e)}"
    
    text = None
    errors = []
    
    # Method 1: Try pdfplumber FIRST and most thoroughly (usually best for Google PDFs)
    if pdfplumber:
        try:
            with pdfplumber.open(file_path) as pdf:
                if not pdf.pages:
                    errors.append("pdfplumber: PDF has no pages")
                else:
                    extracted_text = ""
                    page_count = len(pdf.pages)
                    
                    for page_num, page in enumerate(pdf.pages):
                        try:
                            # Different extraction methods for robustness
                            page_text = ""
                            
                            # Try standard extract_text first
                            try:
                                page_text = page.extract_text() or ""
                            except:
                                pass
                            
                            # If no text, try other methods
                            if not page_text.strip():
                                try:
                                    # Try with layout parameter
                                    page_text = page.extract_text(layout=True) or ""
                                except:
                                    pass
                            
                            # Add page even if empty to preserve document structure
                            extracted_text += f"\n--- Page {page_num + 1} ---\n"
                            extracted_text += page_text
                            
                            # Try to extract tables
                            try:
                                tables = page.extract_tables()
                                if tables:
                                    for table in tables:
                                        for row in table:
                                            extracted_text += " | ".join(str(cell or "") for cell in row) + "\n"
                            except:
                                pass
                        except Exception as e:
                            extracted_text += f"\n--- Page {page_num + 1} ---\n[Unable to extract: {str(e)[:50]}]\n"
                    
                    extracted_text = extracted_text.strip()
                    if extracted_text:
                        text = extracted_text
                    else:
                        errors.append(f"pdfplumber: Extracted text is empty from {page_count} pages")
        except Exception as e:
            errors.append(f"pdfplumber error: {str(e)[:100]}")
    else:
        errors.append("pdfplumber not installed - INSTALL: pip install pdfplumber")
    
    # Method 2: Try pdfminer (different algorithm)
    if (text is None or not text.strip()) and pdfminer_extract_text:
        try:
            extracted_text = pdfminer_extract_text(file_path)
            if extracted_text and extracted_text.strip():
                text = extracted_text
            else:
                errors.append("pdfminer: No text extracted")
        except Exception as e:
            errors.append(f"pdfminer error: {str(e)[:100]}")
    
    # Method 3: Try PyPDF2 (fallback)
    if (text is None or not text.strip()) and PdfReader:
        try:
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                if not reader.pages:
                    errors.append("PyPDF2: PDF has no pages")
                else:
                    extracted_text = ""
                    for page_num, page in enumerate(reader.pages):
                        try:
                            page_text = (page.extract_text() or "").strip()
                            extracted_text += f"\n--- Page {page_num + 1} ---\n"
                            extracted_text += page_text
                        except Exception as e:
                            extracted_text += f"\n--- Page {page_num + 1} ---\n[Error: {str(e)[:50]}]\n"
                    
                    extracted_text = extracted_text.strip()
                    if extracted_text:
                        text = extracted_text
                    else:
                        errors.append("PyPDF2: No text extracted")
        except Exception as e:
            errors.append(f"PyPDF2 error: {str(e)[:100]}")
    else:
        if not PdfReader:
            errors.append("PyPDF2 not installed - INSTALL: pip install PyPDF2")
    
    # Method 4: Try PyMuPDF (fitz) for tougher PDFs
    if (text is None or not text.strip()) and fitz:
        try:
            extracted_text = ""
            with fitz.open(file_path) as doc:
                for page_num, page in enumerate(doc):
                    try:
                        page_text = page.get_text("text") or ""
                        extracted_text += f"\n--- Page {page_num + 1} ---\n"
                        extracted_text += page_text
                    except Exception as e:
                        extracted_text += f"\n--- Page {page_num + 1} ---\n[fitz error: {str(e)[:50]}]\n"
            extracted_text = extracted_text.strip()
            if extracted_text:
                text = extracted_text
            else:
                errors.append("PyMuPDF: No text extracted")
        except Exception as e:
            errors.append(f"PyMuPDF error: {str(e)[:100]}")
    elif (text is None or not text.strip()) and not fitz:
        errors.append("PyMuPDF not installed - INSTALL: pip install PyMuPDF")
    
    # Method 5: Raw byte scan for embedded text strings
    if (text is None or not text.strip()):
        raw_extracted = _extract_text_from_pdf_bytes(file_path)
        if raw_extracted:
            text = raw_extracted
        else:
            errors.append("raw scan: No readable plain text segments found in PDF bytes")
    
    # Method 5: OCR for scanned/image-based PDFs
    if (text is None or not text.strip()) and pytesseract and TESSERACT_CMD:
        try:
            from pdf2image import convert_from_path
            try:
                if POPPLER_CMD:
                    images = convert_from_path(
                        file_path,
                        dpi=300,
                        first_page=1,
                        last_page=5,
                        poppler_path=POPPLER_CMD,
                    )
                else:
                    images = convert_from_path(file_path, dpi=300, first_page=1, last_page=5)
                if images:
                    extracted_text = ""
                    for i, page_image in enumerate(images):
                        try:
                            page_text = pytesseract.image_to_string(page_image, lang='eng')
                            if page_text.strip():
                                extracted_text += f"\n--- Page {i + 1} (OCR) ---\n"
                                extracted_text += page_text
                        except Exception as e:
                            extracted_text += f"\n--- Page {i + 1} ---\n[OCR Error: {str(e)[:50]}]\n"
                    extracted_text = extracted_text.strip()
                    if extracted_text:
                        text = extracted_text
                    else:
                        errors.append("OCR: No text extracted from images")
                else:
                    errors.append("OCR: Could not convert PDF to images")
            except Exception as e:
                errors.append(f"OCR conversion error: {str(e)[:100]}")
        except ImportError:
            errors.append("pdf2image not installed - INSTALL: pip install pdf2image")
        except Exception as e:
            errors.append(f"OCR setup error: {str(e)[:100]}")
    elif (text is None or not text.strip()) and pytesseract and not TESSERACT_CMD:
        errors.append("Tesseract OCR executable not found. Install Tesseract OCR and make sure it is on PATH.")
    else:
        if not pytesseract:
            errors.append("pytesseract not installed - INSTALL: pip install pytesseract")
    
    # Return results
    if text and text.strip():
        return text[:max_chars], None
    
    # All methods failed - provide detailed error
    error_details = "\n".join([f"  • {e}" for e in errors])
    error_msg = f"""Unable to extract text from PDF.
    
Extraction methods tried:
{error_details}

File Details:
  • Path: {file_path}
  • Size: {file_size} bytes
  • Readable: Yes

Troubleshooting:
1. INSTALL missing libraries: pip install pdfplumber pdfminer.six PyPDF2 pytesseract pdf2image
2. Install Tesseract OCR and Poppler for Windows; add Tesseract to PATH.
3. Try opening the PDF in Adobe Reader - ensure it has actual text content
4. If scanned: OCR will attempt to read the pages, but some scanned PDFs need better image quality.
5. For Google Docs PDFs: Export directly from Google Docs as PDF
6. Try uploading a different PDF to test"""
    
    return None, error_msg

def extract_docx_text(file_path, max_chars=50000):
    """Extract text from DOCX files"""
    if not DocxDocument:
        return None, "python-docx not installed. Install with: pip install python-docx"
    
    try:
        doc = DocxDocument(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        for table in doc.tables:
            text += "\n--- TABLE ---\n"
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                text += row_text + "\n"
        return text[:max_chars], None
    except Exception as e:
        return None, f"DOCX extraction error: {str(e)}"

def extract_doc_text(file_path):
    """Extract text from legacy DOC files using docx conversion"""
    # Legacy DOC files require special handling - try to convert or use placeholder
    try:
        # Attempt to open with python-docx (limited support)
        return extract_docx_text(file_path)
    except Exception as e:
        return None, f"DOC extraction error: {str(e)}. Legacy DOC format requires external conversion tools."

def extract_xlsx_text(file_path, max_chars=50000):
    """Extract text from Excel files"""
    if not load_workbook:
        return None, "openpyxl not installed. Install with: pip install openpyxl"
    
    try:
        wb = load_workbook(file_path)
        text = ""
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text += f"\n--- Sheet: {sheet_name} ---\n"
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                text += row_text + "\n"
        return text[:max_chars], None
    except Exception as e:
        return None, f"Excel extraction error: {str(e)}"

def extract_pptx_text(file_path, max_chars=50000):
    """Extract text from PowerPoint files"""
    if not Presentation:
        return None, "python-pptx not installed. Install with: pip install python-pptx"
    
    try:
        prs = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            text += f"\n--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text[:max_chars], None
    except Exception as e:
        return None, f"PowerPoint extraction error: {str(e)}"

def extract_image_text(file_path, max_chars=50000):
    """Extract text/description from image files"""
    try:
        img = Image.open(file_path)
        img_format = img.format or "Unknown"
        width, height = img.size
        
        # For now, return image metadata since OCR requires additional dependencies
        text = f"Image file: {img_format} format\nDimensions: {width}x{height} pixels\nMode: {img.mode}\n"
        text += f"\nFor detailed image analysis, upload and I can describe its contents using AI vision."
        return text[:max_chars], None
    except Exception as e:
        return None, f"Image extraction error: {str(e)}"

def extract_text_file(file_path, max_chars=50000):
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return text[:max_chars], None
    except Exception as e:
        return None, f"Text file extraction error: {str(e)}"

def extract_file_content(file_path, max_chars=50000):
    """
    Universal file extraction function that handles multiple file formats.
    Returns (content, error) tuple.
    """
    if not os.path.exists(file_path):
        return None, f"File not found: {file_path}"
    
    file_ext = os.path.splitext(file_path)[1].lower()
    
    # PDF files
    if file_ext == '.pdf':
        return extract_pdf_text(file_path, max_chars)
    
    # Word documents
    elif file_ext == '.docx':
        return extract_docx_text(file_path, max_chars)
    elif file_ext == '.doc':
        return extract_doc_text(file_path)
    
    # Excel spreadsheets
    elif file_ext in ['.xlsx', '.xls']:
        return extract_xlsx_text(file_path, max_chars)
    
    # PowerPoint presentations
    elif file_ext in ['.pptx', '.ppt']:
        return extract_pptx_text(file_path, max_chars)
    
    # Images
    elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff']:
        return extract_image_text(file_path, max_chars)
    
    # Plain text files
    elif file_ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.yaml', '.yml', '.csv']:
        return extract_text_file(file_path, max_chars)
    
    # For unknown formats, try to read as text
    else:
        return extract_text_file(file_path, max_chars)

#==========================
#SELF-MODIFICATION
#==========================
def modify_self_code(request, auto_apply=False):
    """
    Allow Jarvis to modify its own code 
    Enhanced with auto-apply capability and better safety checks
    """
    try: 
        self_file = __file__ 

        with open(self_file, 'r', encoding='utf-8') as f: 
            current_code = f.read()

        prompt = f"""
    You are Jarvis, an AI assistant that can modify its own code. 

    Current task: {request}

    Your current code file is a Python script. You need to: 
    1. Understand what feature the user wants
    2. Generate the Python code to add that feature
    3. Explain where to add it in the existing code
    4. Ensure the code is safe and won't break existing functionality

    Provide your response in this format: 
    FEATURE: [brief description]
    LOCATION: [where to add - e.g., "after handle_games function" or "in CONFIG section"]
    CODE: 
    ```python
    [the actual Python code to add]
    ```
    EXPLANATION: [how it works]
    SAFETY_CHECK: [any potential risks or considerations]
    
    Keep it simple and safe. Don't remove existing functionality. 
    Only add new functions or modify existing ones carefully.
    """
        response = ask_ai_with_system_prompt( 
            prompt, 
            f"Your name is {ASSISTANT_NAME}. You are a helpful AI that can code and modify yourself. Be careful and thoughtful."
        )

        feature_match = re.search(r'FEATURE:\s*(.+)', response)
        location_match = re.search(r'LOCATION:\s*(.+)', response)
        code_match = re.search(r"CODE:\s*```python\s*(.*?)\s*```", response, re.DOTALL)
        explanation_match = re.search(r'EXPLANATION:\s*(.+)', response, re.DOTALL)
        safety_match = re.search(r'SAFETY_CHECK:\s*(.+)', response, re.DOTALL)

        if not code_match: 
            return f"I understand you want: {request}\n\nBut I need to think more carefully about how to implement this safely. Can you be more specific about what feature you'd like?"

        feature = feature_match.group(1).strip() if feature_match else "New feature"
        location = location_match.group(1).strip() if location_match else "appropriate location"
        new_code = code_match.group(1).strip()
        explanation = explanation_match.group(1).strip() if explanation_match else "Feature added"
        safety_note = safety_match.group(1).strip() if safety_match else "No specific safety concerns"

        backup_file = self_file + ".backup"
        with open(backup_file, 'w', encoding='utf-8') as f: 
            f.write(current_code)

        if auto_apply and "dangerous" not in safety_note.lower() and "risky" not in safety_note.lower():
            try:
                # Find insertion point
                if "after" in location.lower():
                    func_match = re.search(r'after\s+(.+?)\s+function', location.lower())
                    if func_match:
                        func_name = func_match.group(1)
                        pattern = rf'def {func_name}\(.*?\):\s*\n(.*?\n)*?\n'
                        match = re.search(pattern, current_code)
                        if match:
                            insertion_point = match.end()
                            modified_code = current_code[:insertion_point] + f"\n\n# {feature}\n{new_code}\n" + current_code[insertion_point:]
                        else:
                            modified_code = current_code + f"\n\n# {feature}\n{new_code}\n"
                    else:
                        modified_code = current_code + f"\n\n# {feature}\n{new_code}\n"
                else:
                    modified_code = current_code + f"\n\n# {feature}\n{new_code}\n"
                
                # Apply the change
                with open(self_file, 'w', encoding='utf-8') as f:
                    f.write(modified_code)
                
                return f"""
        u2705 SELF-MODIFICATION APPLIED
        u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015

        Feature: {feature}
        Location: {location}
        
        Code has been automatically applied to {self_file}
        
        Explanation: {explanation}
        Safety Note: {safety_note}
        
        Backup created at: {backup_file}
        
        u26a0ufe0f Restart Jarvis to activate the new feature.
        """
            except Exception as e:
                return f"Failed to auto-apply change: {str(e)}\n\nManual application required."

        # Return proposal for manual review
        result = f"""
        u270a SELF-MODIFICATION PROPOSAL
        u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015

        Feature: {feature}
        Location: {location}

        Proposed Code:
        u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015
        {new_code}
        u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015u2015

        Explanation: {explanation}
        Safety Note: {safety_note}

        u26a0ufe0fSAFETY NOTE: 
        This code modification has been proposed but NOT applied automatically. 
        To apply it: 
        1. Review the code above carefully
        2. Copy the code 
        3. Add it to {self_file} at the specified location 
        4. Restart Jarvis 

        A backup has been created at: {backup_file}

        Type "auto apply" to automatically apply this change (if safe).
        """
        
        return result
        
    except Exception as e:
        return f"Error during self-modification: {str(e)}\n\nFor safety, I cannot proceed without understanding the error."


def analyze_self_for_improvements():
    """Autonomous self-analysis to identify improvement areas"""
    try:
        self_file = __file__
        
        with open(self_file, 'r', encoding='utf-8') as f:
            current_code = f.read()
        
        analysis_prompt = f"""
        Analyze this AI assistant code for potential improvements and new features.
        
        Current capabilities: File search, voice control, CAD design, 3D hand tracking, web search, news, music, video, games, app automation
        
        Look for:
        1. Missing common AI assistant features
        2. Inefficiencies in current code
        3. User experience improvements
        4. Safety and error handling gaps
        5. Performance optimizations
        6. Integration opportunities
        
        Provide 3-5 specific, implementable improvements ranked by priority.
        For each improvement, include:
        - PRIORITY: [HIGH/MEDIUM/LOW]
        - FEATURE: [clear description]
        - BENEFIT: [what it improves]
        - IMPLEMENTATION: [brief implementation approach]
        
        Code to analyze:
        ```python
        {current_code[:15000]}
        ``` (truncated for analysis)
        """
        
        analysis = ask_ai_with_system_prompt(
            analysis_prompt,
            "You are an expert AI system architect specializing in self-improving AI assistants. Focus on practical, high-impact improvements."
        )
        
        return analysis
        
    except Exception as e:
        return f"Self-analysis failed: {str(e)}"

def implement_autonomous_improvement(improvement_description):
    """Implement an improvement without user specifying exact code"""
    try:
        self_file = __file__
        
        with open(self_file, 'r', encoding='utf-8') as f:
            current_code = f.read()
        
        implementation_prompt = f"""
        You are Jarvis improving yourself autonomously.
        
        IMPROVEMENT TO IMPLEMENT: {improvement_description}
        
        Generate complete, working Python code to add this improvement to the existing assistant.
        
        Requirements:
        1. Code must be safe and not break existing functionality
        2. Follow existing code patterns and style
        3. Include proper error handling
        4. Add necessary imports if needed
        5. Integrate with existing systems (tone, voice, etc.)
        
        Provide response in this format:
        FEATURE: [brief description]
        LOCATION: [where to add code]
        CODE:
        ```python
        [complete implementation]
        ```
        SAFETY_CHECK: [risk assessment]
        INTEGRATION: [how it connects to existing systems]
        
        Current code structure:
        ```python
        {current_code[:20000]}
        ``` (truncated for context)
        """
        
        response = ask_ai_with_system_prompt(
            implementation_prompt,
            f"Your name is {ASSISTANT_NAME}. You are implementing autonomous self-improvements safely and effectively."
        )
        
        # Parse and apply the improvement
        feature_match = re.search(r'FEATURE:\s*(.+)', response)
        location_match = re.search(r'LOCATION:\s*(.+)', response)
        code_match = re.search(r"CODE:\s*```python\s*(.*?)\s*```", response, re.DOTALL)
        safety_match = re.search(r'SAFETY_CHECK:\s*(.+)', response, re.DOTALL)
        integration_match = re.search(r'INTEGRATION:\s*(.+)', response, re.DOTALL)
        
        if not code_match:
            return f"Could not generate safe implementation for: {improvement_description}"
        
        # Safety checks
        safety_note = safety_match.group(1).strip() if safety_match else "No specific concerns"
        if "dangerous" in safety_note.lower() or "risky" in safety_note.lower():
            return f"Improvement skipped due to safety concerns: {safety_note}"
        
        # Create backup
        backup_file = self_file + f".improvement_backup_{int(time.time())}"
        with open(backup_file, 'w', encoding='utf-8') as f:
            f.write(current_code)
        
        # Apply improvement
        feature = feature_match.group(1).strip() if feature_match else "Autonomous improvement"
        location = location_match.group(1).strip() if location_match else "appropriate location"
        new_code = code_match.group(1).strip()
        integration = integration_match.group(1).strip() if integration_match else "Integrated with existing systems"
        
        # Smart insertion logic
        if "after" in location.lower():
            func_match = re.search(r'after\s+(.+?)\s+function', location.lower())
            if func_match:
                func_name = func_match.group(1)
                pattern = rf'def {func_name}\(.*?\):\s*\n(.*?\n)*?\n'
                match = re.search(pattern, current_code)
                if match:
                    insertion_point = match.end()
                    modified_code = current_code[:insertion_point] + f"\n\n# {feature}\n{new_code}\n" + current_code[insertion_point:]
                else:
                    modified_code = current_code + f"\n\n# {feature}\n{new_code}\n"
            else:
                modified_code = current_code + f"\n\n# {feature}\n{new_code}\n"
        else:
            modified_code = current_code + f"\n\n# {feature}\n{new_code}\n"
        
        # Apply the change
        with open(self_file, 'w', encoding='utf-8') as f:
            f.write(modified_code)
        
        return f"""
🚀 AUTONOMOUS IMPROVEMENT APPLIED

Feature: {feature}
Location: {location}
Integration: {integration}

Code has been automatically applied to {self_file}

Safety Assessment: {safety_note}

Backup created at: {backup_file}

✨ Restart Jarvis to activate new improvement.
        """
        
    except Exception as e:
        return f"Autonomous improvement failed: {str(e)}"

# =========================
# CONTINUOUS SELF-IMPROVEMENT SYSTEM
# =========================
improvement_queue = queue.Queue()
improvement_log = []
continuous_improvement_thread = None

def load_improvement_log():
    """Load improvement history"""
    if os.path.exists(IMPROVEMENT_LOG_FILE):
        try:
            with open(IMPROVEMENT_LOG_FILE, 'r') as f:
                return json.load(f)
        except:
            return []
    return []

def save_improvement_log():
    """Save improvement history"""
    try:
        with open(IMPROVEMENT_LOG_FILE, 'w') as f:
            json.dump(improvement_log, f, indent=2)
    except Exception as e:
        print(f"Failed to save improvement log: {e}")

def log_improvement(improvement_type, description, result):
    """Log an improvement attempt"""
    entry = {
        "timestamp": datetime.now().isoformat(),
        "type": improvement_type,
        "description": description,
        "result": result,
        "success": "APPLIED" in result or "COMPLETED" in result
    }
    improvement_log.append(entry)
    save_improvement_log()

def continuous_improvement_worker():
    """Background worker for continuous self-improvement"""
    global continuous_improvement_thread
    
    while ENABLE_CONTINUOUS_IMPROVEMENT:
        try:
            # Check for queued improvements first
            try:
                improvement = improvement_queue.get(timeout=60)
                if improvement:
                    result = implement_autonomous_improvement(improvement)
                    log_improvement("queued", improvement, result)
                    print(f"🚀 Queued improvement applied: {improvement}")
            except queue.Empty:
                pass
            
            # Periodic self-analysis
            analysis = analyze_self_for_improvements()
            if "Self-analysis failed" not in analysis:
                # Extract high-priority improvements from analysis
                lines = analysis.split('\n')
                for line in lines:
                    if "PRIORITY: HIGH" in line and "FEATURE:" in line:
                        # Extract feature description
                        feature_start = line.find("FEATURE:") + 9
                        if feature_start > 9:
                            feature_end = line.find("BENEFIT:") if "BENEFIT:" in line else len(line)
                            feature = line[feature_start:feature_end].strip()
                            
                            # Check if already implemented
                            already_implemented = any(
                                entry["description"] == feature and entry["success"]
                                for entry in improvement_log[-10:]  # Check last 10 improvements
                            )
                            
                            if not already_implemented:
                                result = implement_autonomous_improvement(feature)
                                log_improvement("autonomous", feature, result)
                                print(f"✨ Autonomous improvement applied: {feature}")
                                break  # Only apply one improvement per cycle
            
            time.sleep(IMPROVEMENT_INTERVAL)
            
        except Exception as e:
            print(f"Continuous improvement error: {e}")
            time.sleep(300)  # Wait 5 minutes before retry

def start_continuous_improvement():
    """Start the continuous improvement background process"""
    global continuous_improvement_thread
    
    if not ENABLE_CONTINUOUS_IMPROVEMENT:
        return
        
    if continuous_improvement_thread is None or not continuous_improvement_thread.is_alive():
        improvement_log.clear()
        improvement_log.extend(load_improvement_log())
        
        continuous_improvement_thread = threading.Thread(
            target=continuous_improvement_worker,
            daemon=True,
            name="ContinuousImprovement"
        )
        continuous_improvement_thread.start()
        print("🔄 Continuous self-improvement system started")

def queue_improvement(improvement_description):
    """Queue an improvement for background processing"""
    if ENABLE_CONTINUOUS_IMPROVEMENT:
        improvement_queue.put(improvement_description)
        return f"Improvement queued: {improvement_description}"
    else:
        return "Continuous improvement is disabled"

def get_improvement_status():
    """Get current improvement system status"""
    if not ENABLE_CONTINUOUS_IMPROVEMENT:
        return "Continuous improvement is disabled"
    
    status = f"""
🔄 Continuous Improvement Status
Active: Yes
Interval: {IMPROVEMENT_INTERVAL} seconds
Queue size: {improvement_queue.qsize()}
Recent improvements: {len([e for e in improvement_log[-5:] if e['success']])} applied
Failed attempts: {len([e for e in improvement_log[-10:] if not e['success']])}
"""
    return status

def get_ai_model_setup_guide():
    """Comprehensive guide for downloading and setting up AI models"""
    return f"""
🤖 AI MODEL SETUP GUIDE

CURRENT MODEL: {AI_MODEL}

📋 WHAT YOU NEED TO DOWNLOAD:

1️⃣ OLLAMA (Required)
   Download from: https://ollama.com/download
   Purpose: Run AI models locally
   Installation: 
   - Windows: Download .exe installer
   - Mac: brew install ollama
   - Linux: curl -fsSL https://ollama.com/install.sh | sh

2️⃣ RECOMMENDED MODELS:

🚀 LLAMA3.1:8B (Current - Fast, Low Internet Dependency)
   Size: 4.9GB
   RAM: 8GB minimum
   Command: ollama run llama3.1:8b
   Context: 128K tokens
   Best for: Strong reasoning, coding, analysis

⚡ QWEN2.5:7B (Faster, Medium Internet Dependency)  
   Size: 4.7GB
   RAM: 8GB minimum
   Command: ollama run qwen2.5:7b
   Context: 32K tokens
   Best for: Speed, recent knowledge, general tasks

🔥 PHI3:MINI (Fastest, Very Low Internet Dependency)
   Size: 2.3GB
   RAM: 4GB minimum
   Command: ollama run phi3:mini
   Context: 4K tokens
   Best for: Extreme speed, minimal resources

📥 INSTALLATION STEPS:

1. INSTALL OLLAMA
   Download from: https://ollama.com/download
   Purpose: Run AI models locally
   Installation: 
   - Windows: Download .exe installer
   - Mac: brew install ollama
   - Linux: curl -fsSL https://ollama.com/install.sh | sh

2. VERIFY INSTALLATION
   Command: ollama --version
   Should show version info

3. DOWNLOAD MODEL (Automatic)
   ollama pull llama3.1:8b    # Current model
   ollama pull qwen2.5:7b     # Faster option
   ollama pull phi3:mini        # Fastest option

4. TEST MODEL
   ollama run llama3.1:8b "Hello, test message"
   Should respond quickly

🔧 CONFIGURATION:
- Models download to: ~/.ollama/models/
- No additional setup required
- Jarvis automatically detects available models
- Restart Jarvis after model changes

💡 RECOMMENDATIONS:
- Start with llama3.1:8b (balanced)
- Upgrade to qwen2.5:7b for more speed
- Use phi3:mini for maximum speed on older hardware

🌐 ALTERNATIVE DOWNLOADS:
- Model Hub: https://ollama.com/library
- GitHub: https://github.com/ollama/ollama
- Documentation: https://github.com/ollama/ollama/blob/main/docs.md

⚠️ SYSTEM REQUIREMENTS:
- RAM: 8GB+ recommended for llama3.1:8b
- Storage: 10GB+ for multiple models
- CPU: Modern multi-core recommended
- GPU: Optional but improves performance

🔄 MODEL MANAGEMENT:
- List models: ollama list
- Remove models: ollama rm <model-name>
- Update models: ollama pull <model-name>:latest

Jarvis will automatically configure and optimize for any downloaded model!
"""

def upgrade_ai_model(new_model=None):
    """Upgrade AI model with validation and configuration updates"""
    global AI_MODEL
    try:
        if not new_model:
            return get_ai_model_setup_guide()
        
        # Validate model is available
        valid_models = ["llama3.1:8b", "qwen2.5:7b", "qwen2.5:1.5b", "mistral:7b", "phi3:mini"]
        
        if new_model not in valid_models:
            return f"❌ Invalid model. Available: {', '.join(valid_models)}\n\n{get_ai_model_setup_guide()}"
        
        # Check if Ollama is installed
        try:
            import subprocess
            result = subprocess.run(['ollama', '--version'], capture_output=True, text=True)
            ollama_installed = result.returncode == 0
        except:
            ollama_installed = False
        
        if not ollama_installed:
            return f"""❌ OLLAMA NOT FOUND

Please install Ollama first:

{get_ai_model_setup_guide()}

After installation, restart Jarvis and try again.
        """
        
        try:
            # Download and setup model
            print(f"Downloading {new_model}...")
            import subprocess
            result = subprocess.run(['ollama', 'pull', new_model], capture_output=True, text=True)
            
            if result.returncode == 0:
                # Create backup of current config
                backup_config = {
                    "old_model": AI_MODEL,
                    "new_model": new_model,
                    "timestamp": datetime.now().isoformat()
                }
                
                with open("model_upgrade_backup.json", "w") as f:
                    json.dump(backup_config, f, indent=2)
                
                # Update model variable
                AI_MODEL = new_model
                
                # Model-specific configurations
                model_configs = {
                    "llama3.1:8b": {
                        "speed": "Fast",
                        "internet_dependency": "Low",
                        "system_prompt": "You are a helpful AI assistant powered by llama3.1:8b with strong local reasoning and 128K context window."
                    },
                    "qwen2.5:7b": {
                        "speed": "Very Fast",
                        "internet_dependency": "Medium",
                        "system_prompt": "You are a helpful AI assistant powered by qwen2.5:7b with excellent recent knowledge and 32K context window."
                    },
                    "phi3:mini": {
                        "speed": "Very Fast",
                        "internet_dependency": "Very Low",
                        "system_prompt": "You are a helpful AI assistant powered by phi3:mini with extreme speed and 4K context window."
                    }
                }
                
                config = model_configs.get(new_model, model_configs["llama3.1:8b"])
                
                return f"""
🚀 AI MODEL UPGRADED SUCCESSFULLY

Previous model: {backup_config['old_model']}
New model: {new_model}
Speed: {config['speed']}
Internet dependency: {config['internet_dependency']}
Context window: {config.get('context', '128K')}

✨ Model downloaded and configured!
Restart Jarvis to activate new model.
Backup created: model_upgrade_backup.json

🔄 OLLAMA COMMANDS FOR REFERENCE:
- ollama list                    # Show all models
- ollama pull {model}           # Download new model
- ollama run {model}             # Start model
- ollama rm {model}             # Remove model
                """
            else:
                return f"❌ Failed to download {new_model}: {result.stderr}"
        
        except Exception as e:
            return f"❌ Model upgrade failed: {str(e)}\n\nFor manual setup:\n{get_ai_model_setup_guide()}"

    except Exception as e:
        return f"❌ Model upgrade encountered an error: {str(e)}"

def is_self_modification_query(c):
    """Check if user wants Jarvis to modify itself"""
    mod_phrases = [
        "modify yourself",
        "add a feature",
        "add feature",
        "improve yourself",
        "update yourself",
        "change your code",
        "add to your code",
        "give yourself",
        "add the ability",
        "make yourself",
        "enhance yourself",
        "upgrade yourself",
        "improve autonomously",
        "self improve",
        "analyze yourself",
        "make yourself better",
        "autonomous improvement",
        "improvement status",
        "continuous improvement",
        "upgrade model",
        "change model",
        "recommend model",
        "ai model",
        "model upgrade",
    ]
    return contains_any(c, mod_phrases)

def is_voice_control_query(c):
    """Check if user wants to control voice settings"""
    voice_phrases = [
        "mute voice",
        "unmute voice",
        "toggle voice",
        "voice off",
        "voice on",
        "stop talking",
        "start talking",
        "speak louder",
        "speak quieter",
        "speak faster",
        "speak slower",
    ]
    return contains_any(c, voice_phrases)

def is_tone_control_query(c):
    """Check if user wants to control tone settings"""
    tone_phrases = [
        "change tone",
        "set tone",
        "tone to",
        "be more",
        "be less",
        "sound more",
        "sound less",
        "personality",
        "speak like",
        "talk like",
        "tone settings",
        "list tones",
        "available tones",
        "voice",
        "text",
        "format",
        "switch to voice",
        "switch to text",
        "talk to me",
        "speak to me",
    ]
    return contains_any(c, tone_phrases)

# =========================
# APPS
# =========================
BUILT_IN_APPS = {
    "calculator": "calc.exe",
    "notepad": "notepad.exe",
    "paint": "mspaint.exe",
    "cmd": "cmd.exe",
    "steam": r"C:\Program Files (x86)\Steam\Steam.exe",
    "discord": r"%LOCALAPPDATA%\\Discord\\Update.exe --processStart Discord.exe",
}

OPEN_VERBS = ["open", "launch", "start", "run"]
SEARCH_VERBS = ["search", "google", "look up"]
CURRENT_EVENT_PHRASES = [
    "current information",
    "current events",
    "latest news",
    "latest information",
    "recent news",
    "recent updates",
    "breaking news",
    "what is happening",
    "what's happening",
    "what is going on",
    "what's going on",
    "today's news",
]
CURRENT_EVENT_TOPICS = [
    "war",
    "tensions",
    "conflict",
    "politics",
    "election",
    "russia",
    "ukraine",
    "israel",
    "gaza",
    "china",
    "taiwan",
    "news",
    "international relations",
]

# Vision-related triggers
VISION_TRIGGERS = [
    "what do you see",
    "look at me",
    "can you see me",
    "what am i doing",
    "what's in front of you",
    "describe what you see",
    "take a look",
    "what am i wearing",
    "what color is",
    "how many",
    "identify this",
    "recognize this",
    "am i",
    "do i look",
]

# =========================
# ROBLOX
# =========================
def find_roblox():
    global _ROBLOX_CACHE

    if _ROBLOX_CACHE and os.path.exists(_ROBLOX_CACHE):
        return _ROBLOX_CACHE

    base = os.path.join(os.environ.get("LOCALAPPDATA", ""), "Roblox", "Versions")
    if not os.path.exists(base):
        return None

    for folder in os.listdir(base):
        exe = os.path.join(base, folder, "RobloxPlayerBeta.exe")
        if os.path.exists(exe):
            _ROBLOX_CACHE = exe
            return exe

    return None

# =========================
# CAMERA / VISION
# =========================
def capture_image():
    """Capture an image from the webcam"""
    if not ENABLE_CAMERA:
        return None, "Camera is disabled in settings"
    
    try:
        cap = cv2.VideoCapture(CAMERA_INDEX)
        if not cap.isOpened():
            return None, "Could not access camera"
        
        # Let camera warm up
        for _ in range(5):
            cap.read()
        
        ret, frame = cap.read()
        cap.release()
        
        if not ret:
            return None, "Failed to capture image"
        
        # Convert BGR to RGB
        frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        
        # Convert to PIL Image
        pil_image = Image.fromarray(frame_rgb)
        
        # Convert to base64
        buffered = io.BytesIO()
        pil_image.save(buffered, format="JPEG", quality=85)
        img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
        
        return img_base64, None
        
    except Exception as e:
        return None, f"Camera error: {str(e)}"

def ask_ai_maybe_with_camera(cmd, use_vision):
    """When use_vision is True (browser camera toggle), attach a fresh webcam frame for chat-style requests."""
    if not use_vision or not ENABLE_CAMERA:
        return ask_ai(cmd)
    img_base64, error = capture_image()
    if error:
        return f"{error}\n\nAnswering without a live image:\n\n{ask_ai(cmd)}"
    return ask_ai_with_vision(cmd, img_base64)

def ask_ai_with_vision(prompt, image_base64=None):
    """Ask AI with optional image input"""
    system_prompt = (
        f"Your name is {ASSISTANT_NAME}. "
        f"You are {ASSISTANT_NAME}, a helpful AI assistant with vision capabilities. "
        f"Never mention your model name (Qwen, Llama, etc). "
        f"If asked your name, always answer that your name is {ASSISTANT_NAME}. "
        "When describing images, be concise and focus on what the user asks about. "
        "Keep responses brief and helpful."
    )
    
    try:
        if image_base64:
            # Use vision model
            payload = {
                "model": VISION_MODEL,
                "messages": [
                    {
                        "role": "system",
                        "content": system_prompt,
                    },
                    {
                        "role": "user",
                        "content": prompt,
                        "images": [image_base64]
                    }
                ],
                "stream": False
            }
        else:
            # Use regular model
            payload = {
                "model": AI_MODEL,
                "messages": [
                    {
                        "role": "system",
                        "content": system_prompt,
                    },
                    {"role": "user", "content": prompt}
                ],
                "stream": False
            }

        r = http.post(OLLAMA_URL, json=payload, timeout=AI_TIMEOUT)
        r.raise_for_status()
        return r.json()["message"]["content"]

    except ReadTimeout:
        try:
            r = http.post(OLLAMA_URL, json=payload, timeout=AI_RETRY_TIMEOUT)
            r.raise_for_status()
            return r.json()["message"]["content"]
        except ReadTimeout:
            return (
                f"{ASSISTANT_NAME} is taking too long to respond. "
                f"Your model may be too slow for this request. "
                f"Try again, use a smaller/faster model, or increase the timeout."
            )
            return f"AI error: {str(e)}"

        return f"AI error: {str(e)}"

# =========================
# UTILS
# =========================
def normalize(t):
    return t.lower().strip()

def score(query, keywords):
    return max(fuzz.partial_ratio(query, k) for k in keywords)

def fuzzy_contains(text, phrases, threshold=80):
    return any(fuzz.partial_ratio(text, phrase) >= threshold for phrase in phrases)

def contains_any(text, phrases):
    return any(phrase in text for phrase in phrases)

def strip_leading_action(text, verbs):
    cleaned = normalize(text)
    for verb in verbs:
        if cleaned.startswith(verb + " "):
            return cleaned[len(verb):].strip()
    return cleaned

def starts_with_fuzzy_verb(text, verbs, threshold=85):
    words = normalize(text).split()
    if not words:
        return False
    first = words[0]
    return any(fuzz.ratio(first, verb) >= threshold for verb in verbs)

def is_greeting(c):
    return any(x in c for x in ["hello", "hi", "hey", "yo", "jarvis", "wake up"])

def is_vision_query(c):
    """Check if user wants to use camera/vision"""
    return contains_any(c, VISION_TRIGGERS)

def is_action_command(c):
    """Check if this is clearly an action command (not a question)"""
    # Strong action indicators
    action_starters = OPEN_VERBS + SEARCH_VERBS + [
        "play", "launch", "start", "run", "open", "search", "google", 
        "listen", "watch", "show me"
    ]
    
    # Check if starts with action verb
    words = c.split()
    if words and words[0] in action_starters:
        return True
    
    # Check for specific action patterns
    action_patterns = [
        r"^(open|launch|start|run)\s+\w+",
        r"^(play|listen to|watch)\s+",
        r"^(search|google|look up)\s+",
        r"^(show me|find me)\s+\w+\s+(on|in)\s+(youtube|netflix|spotify)",
    ]
    
    return any(re.match(pattern, c) for pattern in action_patterns)

def is_question(c):
    """Improved question detection that distinguishes from action commands"""
    # First check if it's clearly an action command
    if is_action_command(c):
        return False
    
    # Question words that indicate information seeking
    question_words = [
        "what", "who", "why", "how", "when", "where", "which", 
        "explain", "define", "tell me about", "talk to me about", 
        "discuss", "describe", "is", "are", "was", "were", 
        "do", "does", "did", "can", "could", "would", "should", "write"
    ]
    
    # Has question mark
    if "?" in c:
        return True
    
    # Starts with question word (but not followed by action like "what time is it" -> time intent)
    words = c.split()
    if words and words[0] in question_words:
        # Exclude time queries
        if not is_time_query(c):
            return True
    
    # Conversational question patterns
    conversational_patterns = [
        r"^(tell|explain|describe|discuss)\s+",
        r"^(what|who|why|how|when|where)\s+",
        r"^(can you|could you|would you)\s+",
        r"\b(meaning|definition|explanation)\b",
    ]
    
    return any(re.search(pattern, c) for pattern in conversational_patterns)

def is_time_query(c):
    exact_time_phrases = [
        "time",
        "date",
        "clock",
        "what time is it",
        "what is the time",
        "tell me the time",
        "current time",
        "what's the time",
        "what is today's date",
        "today's date",
    ]
    if c in exact_time_phrases:
        return True

    phrase_patterns = [
        r"\bwhat time is it\b",
        r"\bwhat is the time\b",
        r"\btell me the time\b",
        r"\bcurrent time\b",
        r"\bwhat's the time\b",
        r"\bwhat is today's date\b",
        r"\btoday's date\b",
        r"\bcurrent date\b",
        r"\bwhat is the date\b",
    ]

    return any(re.search(pattern, c) for pattern in phrase_patterns)

def is_watch_query(c):
    """Improved watch detection that distinguishes from questions"""
    # Exclude if it's a pure question about watching
    if is_question(c) and not any(word in c for word in ["youtube", "netflix", "video on"]):
        return False
    
    watch_phrases = [
        "watch",
        "watch something",
        "i want to watch",
        "i wanna watch",
        "lets watch",
        "let's watch",
        "youtube",
        "netflix",
        "stream",
        "movie",
        "video",
        "play video",
        "show me video",
    ]
    return contains_any(c, watch_phrases)

def is_listen_query(c):
    """Improved listen detection"""
    # Exclude pure questions
    if is_question(c) and not any(word in c for word in ["spotify", "play music"]):
        return False
        
    listen_phrases = [
        "listen to something",
        "i want to listen",
        "i wanna listen",
        "listen to music",
        "play music",
        "spotify",
        "song",
        "podcast",
        "play song",
        "play spotify",
    ]
    return contains_any(c, listen_phrases)

def is_game_query(c):
    """Improved game detection"""
    # Exclude pure questions about games
    if is_question(c) and "play" not in c:
        return False
        
    game_phrases = [
        "play",
        "play game",
        "play games",
        "i want to play",
        "i wanna play",
        "i want 2 play",
        "lets play",
        "let's play",
        "game",
        "games",
        "gaming",
        "roblox",
        "open roblox",
        "launch game",
    ]
    return contains_any(c, game_phrases)

def is_current_events_query(c):
    has_current_phrase = contains_any(c, CURRENT_EVENT_PHRASES)
    has_topic = contains_any(c, CURRENT_EVENT_TOPICS)
    asks_for_updates = contains_any(c, ["current", "latest", "recent", "today", "news", "update"])
    return has_current_phrase or (asks_for_updates and has_topic)

def is_cad_model_query(c):
    """True when the user wants a 3D / CAD model (OpenSCAD), not generic create tasks."""
    c = normalize(c)
    if is_question(c):
        purely_info = (
            c.startswith("what is ")
            or c.startswith("what are ")
            or "explain " in c
            or "define " in c
            or c.startswith("why ")
        )
        design_intent = contains_any(
            c,
            [
                "create",
                "design",
                "make",
                "generate",
                "model a",
                "build a",
                "can you design",
                "can you create",
                "can you make",
                "help me design",
                "i need a",
                "i want a",
            ],
        ) or re.search(r"\b(draw|sketch)\b.*\b(3d|cad|openscad)\b", c)
        if purely_info and not design_intent:
            return False

    strong = [
        "openscad",
        ".scad",
        "open scad",
        "scad file",
        "scad code",
        "3d model",
        "3d print",
        "3d design",
        "3d printable",
        "cad model",
        "cad design",
        "parametric model",
        "for 3d printer",
        "3d printer",
        "generate stl",
        "export to stl",
    ]
    if contains_any(c, strong):
        return True

    exclude = [
        "website",
        "web app",
        "webpage",
        "playlist",
        "business plan",
        "essay",
        "story",
        "song ",
        "music video",
        "graphic design",
        "ui design",
        "logo for",
    ]
    if contains_any(c, exclude) and not contains_any(
        c, ["3d", "openscad", "cad", "print", "stl", "scad", "bracket", "stand", "holder", "enclosure"]
    ):
        return False

    phrase_hits = [
        "create a 3d",
        "create a model",
        "create a cad",
        "create an openscad",
        "design a 3d",
        "design a bracket",
        "design a stand",
        "design a holder",
        "design a mount",
        "design a case",
        "design an enclosure",
        "make a 3d",
        "make a bracket",
        "make a stand",
        "make a holder",
        "generate a 3d",
        "generate a cad",
        "model a ",
        "model for 3d",
        "build a 3d",
    ]
    if contains_any(c, phrase_hits):
        return True

    if re.search(
        r"\b(create|design|make|generate)\s+(a|an|the)?\s*(3d|openscad|cad|parametric|printable)\b",
        c,
    ):
        return True
    if re.search(
        r"\b(create|design|make)\s+(a|an)\s+\w+\s+(bracket|stand|holder|enclosure|mount|case|gear|knob|clamp|vase|box)\b",
        c,
    ):
        return True
    return False

def clean_news_query(cmd):
    query = normalize(cmd)

    fillers = [
        "give me",
        "show me",
        "tell me",
        "find me",
        "search",
        "google",
        "look up",
        "news about",
        "news on",
        "latest news about",
        "latest news on",
        "latest information about",
        "latest information on",
        "current information about",
        "current information on",
        "current information regarding",
        "current information between",
        "information about",
        "information on",
        "regarding",
        "about",
    ]

    changed = True
    while changed:
        changed = False
        for filler in fillers:
            if query.startswith(filler + " "):
                query = query[len(filler):].strip()
                changed = True

    query = query.replace("tensions between", "").replace("tensions regarding", "").strip()
    return query or None

def force_intent(cmd, c):
    """Force specific intents for clear action commands"""
    # Vision intent
    if is_vision_query(c):
        return "vision"
    
    # Watch intents
    if contains_any(c, ["watch something", "i wanna watch", "i want to watch", "lets watch"]):
        return "watch"
    if fuzzy_contains(c, ["watch something", "i wanna watch", "i want to watch", "lets watch"], threshold=90):
        return "watch"

    # Listen intents
    if contains_any(c, [
        "i want to listen",
        "i wanna listen",
        "listen to something",
        "i want to listen to something",
        "i wanna listen to something",
        "listen music",
        "listen to music",
        "play music",
        "open spotify",
        "i want music",
    ]):
        return "listen"
    if fuzzy_contains(c, [
        "i want to listen",
        "i wanna listen",
        "listen to something",
        "i want to listen to something",
        "i wanna listen to something",
        "listen music",
        "listen to music",
        "play music",
        "open spotify",
        "i want music",
    ], threshold=85):
        return "listen"

    # Game intents
    if contains_any(c, [
        "i wanna play",
        "i want to play",
        "i want 2 play",
        "wanna play",
        "want to play",
        "play game",
        "play games",
        "lets play",
        "open roblox",
    ]):
        return "game"
    if fuzzy_contains(c, [
        "i wanna play",
        "i want to play",
        "i want 2 play",
        "wanna play",
        "want to play",
        "play game",
        "play games",
        "lets play",
    ], threshold=85):
        return "game"

    # Do not treat OpenSCAD / CAD requests as "open <app>"
    if is_cad_model_query(c):
        return None

    # Action verbs
    if starts_with_fuzzy_verb(c, OPEN_VERBS):
        return "open"

    if starts_with_fuzzy_verb(c, SEARCH_VERBS):
        return "search"

    return None

# =========================
# INTENT ENGINE (IMPROVED)
# =========================
def detect_intent(cmd):
    c = normalize(cmd)

    # Priority 1: Greetings
    if is_greeting(c):
        return "chat", 100, cmd

    if is_self_modification_query(c):
        return "self_modify", 100, cmd

    # Priority 2: Vision queries
    if is_vision_query(c):
        return "vision", 100, cmd

    # Priority 3: Time queries
    if is_time_query(c):
        return "time", 100, cmd

    # Priority 4: News/current events
    if is_current_events_query(c):
        return "news", 100, cmd

    # Priority 4.5: CAD before "open" fuzzy matches ("openscad" contains "open")
    if is_cad_model_query(c):
        return "cad", 100, cmd

    # Priority 5: Check for forced intents (explicit action commands)
    forced = force_intent(cmd, c)
    if forced:
        return forced, 100, cmd

    # Priority 6: Specific media/app queries (before general questions)
    if is_listen_query(c):
        return "listen", 95, cmd

    if is_game_query(c):
        return "game", 95, cmd

    if is_watch_query(c):
        return "watch", 95, cmd

    # Priority 7: Questions go to chat (after checking for action intents)
    if is_question(c):
        return "chat", 95, cmd

    # Priority 8: News mentions
    if contains_any(c, [" news", "latest news"]) or c == "news":
        return "news", 90, cmd

    # Priority 9: Fuzzy matching for remaining commands
    intents = {}
    intents["open"] = 0 if is_cad_model_query(c) else score(c, ["open", "launch", "start", "run"])
    intents["music"] = score(c, ["spotify", "play music"])
    intents["watch"] = score(c, ["watch", "youtube", "netflix", "stream", "movie", "video"])
    intents["listen"] = score(c, ["listen", "podcast", "music", "song", "spotify"])
    intents["news"] = score(c, ["news", "latest"])
    intents["game"] = score(c, ["game", "games", "gaming", "play game", "play games"])
    intents["search"] = score(c, ["search", "google", "look up", "find"])

    best = max(intents, key=intents.get)
    confidence = intents[best]

    thresholds = {
        "music": 80,
        "watch": 80,
        "listen": 80,
        "game": 75,
        "open": 70,
        "news": 60,
        "search": 70,
    }

    if confidence < thresholds.get(best, 999):
        return "chat", 90, cmd

    return best, confidence, cmd

# =========================
# AI
# =========================
def ask_ai(prompt):
    # Enhanced system prompt for the configured model
    system_prompt = get_tone_system_prompt() if ENABLE_TONE_SYSTEM else (
        f"""Your name is {ASSISTANT_NAME}. You are a helpful AI assistant powered by {AI_MODEL}.

Key capabilities:
- Fast response times (under 2 seconds)
- Strong local knowledge and reasoning
- Efficient problem-solving without internet reliance
- Good at coding, analysis, and creative tasks
- Optimized for practical assistance

Guidelines:
- Prioritize speed and accuracy
- Use your built-in knowledge first, only use web search when specifically requested
- Be concise but thorough
- Excel at coding, technical tasks, and problem-solving
- Provide practical solutions over theoretical ones
- Maintain context and remember previous interactions

Current model: {AI_MODEL}"""
    )
    return ask_ai_with_system_prompt(prompt, system_prompt, model=AI_MODEL)

def ask_ai_with_system_prompt(prompt, system_prompt, model=None, timeout=None, retry_timeout=None):
    model = model or AI_MODEL
    timeout = timeout or AI_TIMEOUT
    retry_timeout = retry_timeout or AI_RETRY_TIMEOUT
    try:
        payload = {
            "model": model,
            "messages": [
                {
                    "role": "system",
                    "content": system_prompt,
                },
                {"role": "user", "content": prompt}
            ],
            "stream": False
        }

        r = http.post(OLLAMA_URL, json=payload, timeout=timeout)
        r.raise_for_status()
        return r.json()["message"]["content"]

    except ReadTimeout:
        try:
            r = http.post(OLLAMA_URL, json=payload, timeout=retry_timeout)
            r.raise_for_status()
            return r.json()["message"]["content"]
        except ReadTimeout:
            fallback_model = None
            
            # Smart fallback: for documents, prefer capable small models
            if model == "qwen2.5:7b":
                fallback_model = "qwen2.5:1.5b"
            elif model == "qwen2.5:1.5b":
                fallback_model = "phi3:mini"
            elif model != "phi3:mini":
                # For other models, try qwen2.5:7b
                fallback_model = "qwen2.5:7b"

            if fallback_model:
                try:
                    payload["model"] = fallback_model
                    r = http.post(OLLAMA_URL, json=payload, timeout=retry_timeout)
                    r.raise_for_status()
                    return r.json()["message"]["content"]
                except Exception:
                    pass

            return (
                f"{ASSISTANT_NAME} is taking too long to respond from Ollama. "
                f"This request may be too large or too complex for the current timeout. "
                f"Try again with a shorter prompt, a faster model, or increase the timeout."
            )

        except Exception as e:
            return f"AI error: {str(e)}"

    except Exception as e:
        return f"AI error: {str(e)}"

def review_file_content(filename, content_or_path, user_request="", is_file_path=False):
    """
    Review file content from either direct content string or file path.
    Supports PDFs, DOCX, XLSX, PPTX, images, and text files.
    
    Args:
        filename: Name of the file
        content_or_path: Either file content as string OR file path
        user_request: Specific question/request about the file
        is_file_path: If True, content_or_path is treated as a file path
    """
    safe_name = filename or "uploaded file"
    
    # Extract content based on whether it's a path or direct content
    if is_file_path:
        extracted_content, error = extract_file_content(content_or_path)
        if error:
            google_pdf_hint = ""
            if "pdf" in filename.lower():
                google_pdf_hint = "\n\n📄 For Google PDFs: Google Docs PDFs sometimes use special encoding. Try:\n- Download the PDF from Google Drive again\n- Export as PDF from Google Docs (File > Download > PDF)\n- Use Chrome's Print to PDF feature on the Google Docs page"
            return f"❌ Error extracting file: {error}{google_pdf_hint}"
        if not extracted_content or not extracted_content.strip():
            pdf_guidance = ""
            if ".pdf" in filename.lower():
                pdf_guidance = "\n\n📄 **For PDF files:**\n- Google PDFs: Try exporting from Google Docs directly\n- Scanned PDFs: Use OCR or convert image-based PDFs\n- If your PDF is text-based: Try re-saving it with a different PDF tool\n- Check: Open the PDF in Adobe Reader to verify it has text"
            return f"❌ No text content found in: {safe_name}{pdf_guidance}"
        trimmed_content = extracted_content[:50000]
    else:
        # Direct content string (from web upload)
        trimmed_content = content_or_path[:50000] if content_or_path else ""
    
    # Build the analysis prompt based on file type
    file_ext = os.path.splitext(safe_name)[1].lower()
    
    if user_request:
        prompt = (
            f"You are analyzing a file named '{safe_name}' ({file_ext} format).\n"
            f"User request: \"{user_request}\"\n\n"
            f"Please respond directly to the user's request based on the file contents. "
            "Be clear and concise. Provide direct answers without unnecessary filler.\n\n"
            f"File contents:\n```\n{trimmed_content}\n```"
        )
    else:
        prompt = (
            f"Review this file named '{safe_name}' ({file_ext} format).\n\n"
            f"User request: Check for mistakes, issues, and suggest improvements.\n\n"
            "Give a practical review with these sections:\n"
            "1. Summary\n"
            "2. Key findings (mistakes, risks, issues)\n"
            "3. Suggested improvements\n"
            "4. Improved snippet or example if useful\n\n"
            "Be specific and concise. If there are no major issues, say that clearly.\n\n"
            f"File contents:\n```\n{trimmed_content}\n```"
        )
    
    system_prompt = (
        f"Your name is {ASSISTANT_NAME}. "
        "You are a careful reviewer for code, text, documents, spreadsheets, and presentations. "
        "Focus on mistakes, bugs, clarity issues, maintainability, and practical improvements. "
        "Adapt your review style to the file type. "
        "Do not invent file contents that are not present."
    )
    
    return ask_ai_with_system_prompt(
        prompt,
        system_prompt,
        model=AI_DOC_MODEL,
        timeout=AI_DOC_TIMEOUT,
        retry_timeout=AI_DOC_RETRY_TIMEOUT,
    )

def format_cad_reply_short(design):
    title = str(design.get("title", "CAD design")).strip() or "CAD design"
    summary = str(design.get("summary", "")).strip()
    return f"CAD design ready: {title}. {summary}".strip()


def format_cad_reply_cli(design):
    title = str(design.get("title", "CAD design")).strip() or "CAD design"
    summary = str(design.get("summary", "")).strip()
    fn = str(design.get("filename", "design.scad")).strip() or "design.scad"
    code = str(design.get("code", ""))
    preview = code[:2500] + ("..." if len(code) > 2500 else "")
    return (
        f"=== CAD assistant ({CAD_MODEL}) ===\n{title}\n{summary}\n\n"
        f"File: {fn}\n\n--- OpenSCAD ---\n{preview}"
    )


def generate_cad_design(request, style="openscad"):
    trimmed_request = (request or "").strip()[:4000]
    if not trimmed_request:
        return {
            "title": "Untitled design",
            "summary": "No CAD request was provided.",
            "code": "",
            "filename": "design.scad",
        }

    system_prompt = (
        f"Your name is {ASSISTANT_NAME}. "
        "You are a specialist CAD assistant embedded in Jarvis. "
        "You output OpenSCAD (.scad) source code only inside JSON. "
        "Use millimeters as the default unit; set $fn for smooth cylinders/spheres when needed. "
        "Prefer union(), difference(), hull(), linear_extrude(), rotate_extrude() as appropriate. "
        "Keep geometry simple and printable; avoid infinite recursion or invalid CSG. "
        "Return valid JSON only with keys: title, summary, code, filename. "
        "The code value must be raw OpenSCAD text (escape newlines as \\n in JSON). "
        "Filename must end in .scad."
    )
    prompt = (
        "Create a 3D model from this request.\n\n"
        f"Request: {trimmed_request}\n\n"
        "Return JSON only. "
        "title: short product name. "
        "summary: 1-2 sentences on what you built and print notes if any. "
        "code: complete runnable OpenSCAD. "
        "filename: descriptive_snake_case.scad"
    )

    raw = ask_ai_with_system_prompt(
        prompt,
        system_prompt,
        model=CAD_MODEL,
        timeout=CAD_TIMEOUT,
        retry_timeout=CAD_RETRY_TIMEOUT,
    )

    try:
        start = raw.find("{")
        end = raw.rfind("}")
        if start != -1 and end != -1:
            raw = raw[start:end + 1]
        data = json.loads(raw)
    except Exception:
        data = {
            "title": "CAD concept",
            "summary": "Jarvis generated a design response, but it was not in clean JSON format.",
            "code": raw,
            "filename": "design.scad",
        }

    data["title"] = str(data.get("title", "CAD concept")).strip() or "CAD concept"
    data["summary"] = str(data.get("summary", "")).strip()
    data["code"] = str(data.get("code", "")).strip()
    data["filename"] = str(data.get("filename", "design.scad")).strip() or "design.scad"

    if not data["filename"].lower().endswith(".scad"):
        data["filename"] += ".scad"

    return data


def find_openscad_executable():
    for name in ("openscad", "openscad.com", "openscad.exe"):
        p = shutil.which(name)
        if p:
            return p
    if os.name == "nt":
        for base in (
            os.environ.get("ProgramFiles", r"C:\Program Files"),
            os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)"),
        ):
            for sub in ("OpenSCAD", "OpenSCAD (64 bit)"):
                for exe in ("openscad.exe", "openscad.com"):
                    cand = Path(base) / sub / exe
                    if cand.is_file():
                        return str(cand)
    return None


def compile_scad_to_stl(scad_source):
    """Compile OpenSCAD source to STL bytes. Returns (stl_bytes_or_none, error_or_none)."""
    exe = find_openscad_executable()
    if not exe:
        return None, (
            "OpenSCAD is not installed or not on PATH. "
            "Install from https://openscad.org/ to enable STL preview."
        )
    if not (scad_source or "").strip():
        return None, "No OpenSCAD source to compile."
    try:
        with tempfile.TemporaryDirectory() as tmp:
            tdir = Path(tmp)
            scad_path = tdir / "model.scad"
            stl_path = tdir / "model.stl"
            scad_path.write_text(scad_source, encoding="utf-8")
            r = subprocess.run(
                [exe, "-o", str(stl_path), str(scad_path)],
                capture_output=True,
                text=True,
                timeout=OPENSCAD_COMPILE_TIMEOUT,
            )
            if r.returncode != 0:
                err = (r.stderr or r.stdout or "compile failed").strip()[:800]
                return None, f"OpenSCAD: {err}"
            if not stl_path.exists():
                return None, "OpenSCAD did not produce an STL file."
            return stl_path.read_bytes(), None
    except subprocess.TimeoutExpired:
        return None, "OpenSCAD compile timed out."
        return None, str(e)


def _prune_stl_preview_cache():
    now = time.time()
    dead = [k for k, v in STL_PREVIEW_CACHE.items() if now - v.get("ts", 0) > STL_CACHE_TTL_SEC]
    for k in dead:
        STL_PREVIEW_CACHE.pop(k, None)
    while len(STL_PREVIEW_CACHE) > STL_CACHE_MAX_ITEMS:
        oldest = min(STL_PREVIEW_CACHE.items(), key=lambda kv: kv[1].get("ts", 0))[0]
        STL_PREVIEW_CACHE.pop(oldest, None)


def register_stl_preview(stl_bytes):
    _prune_stl_preview_cache()
    pid = str(uuid.uuid4())
    STL_PREVIEW_CACHE[pid] = {"bytes": stl_bytes, "ts": time.time()}
    return pid


def get_stl_preview_bytes(preview_id):
    _prune_stl_preview_cache()
    ent = STL_PREVIEW_CACHE.get(preview_id)
    if not ent:
        return None
    return ent.get("bytes")


def attach_stl_preview(design):
    """Add stl_preview_id / stl_ready / stl_error for web UI."""
    d = dict(design)
    code = d.get("code") or ""
    stl_bytes, err = compile_scad_to_stl(code)
    if stl_bytes:
        d["stl_preview_id"] = register_stl_preview(stl_bytes)
        d["stl_ready"] = True
        d.pop("stl_error", None)
    else:
        d["stl_ready"] = False
        d["stl_error"] = err or "Could not build STL preview."
    base = (d.get("filename") or "design.scad").replace(".scad", ".stl")
    d["stl_filename"] = base if base.lower().endswith(".stl") else base + ".stl"
    return d


# =========================
# LIVE NEWS FEATURE
# =========================
def fetch_news(query=None):
    try:
        with DDGS() as ddgs:
            results = ddgs.news(
                query=query or "latest news",
                region="wt-wt",
                max_results=5
            )

        if not results:
            return "No news found."

        # Extract news information without URLs
        news_info = []
        for r in results:
            news_info.append(f"Title: {r.get('title', 'No title')}\nSource: {r.get('source', 'Unknown')}\nDate: {r.get('date', 'Unknown')}\nContent: {r.get('body', 'No content available')}")
        
        # Ask AI to summarize the news
        summary_prompt = f"""Please provide a comprehensive summary of the following news articles about "{query or 'latest news'}":

{chr(10).join(news_info)}

Focus on the most important news stories and key information. Do not include URLs or links. Provide a clear, informative summary that highlights the main points."""
        
        summary = ask_ai(summary_prompt)
        return f"Here's the latest news about {query or 'current events'}:\n\n{summary}"

    except Exception as e:
        return f"News error: {str(e)}"

def handle_news(cmd):
    query = clean_news_query(cmd.replace("news", " ").strip())
    if not query:
        query = None

    return fetch_news(query)

# =========================
# HANDLERS
# =========================
def handle_vision(cmd):
    """Handle vision/camera requests"""
    img_base64, error = capture_image()
    
    if error:
        return error
    
    # Extract the specific question from the command
    vision_prompt = cmd
    for trigger in VISION_TRIGGERS:
        if trigger in vision_prompt.lower():
            vision_prompt = vision_prompt.lower().replace(trigger, "").strip()
    
    # If no specific question, use a default
    if not vision_prompt or len(vision_prompt) < 3:
        vision_prompt = "Describe what you see in this image."
    
    return ask_ai_with_vision(vision_prompt, img_base64)

def handle_music():
    webbrowser.open("https://open.spotify.com")
    return "Opening Spotify"

def handle_listen():
    webbrowser.open("https://open.spotify.com")
    return "Opening Spotify"

def handle_watch(cmd, interactive=True):
    cleaned_cmd = strip_leading_action(cmd, OPEN_VERBS + SEARCH_VERBS)

    if "youtube" in cleaned_cmd:
        webbrowser.open("https://youtube.com/results?search_query=" + quote(cleaned_cmd))
        return "Opening YouTube"

    if "netflix" in cleaned_cmd:
        webbrowser.open("https://netflix.com/search?q=" + quote(cleaned_cmd))
        return "Opening Netflix"

    if not interactive:
        return "YouTube or Netflix?"

    choice = input("YouTube or Netflix? >>> ").lower()

    if "youtube" in choice:
        webbrowser.open("https://youtube.com/results?search_query=" + quote(cmd))
        return "Opening YouTube"

    if "netflix" in choice:
        webbrowser.open("https://netflix.com/search?q=" + quote(cmd))
        return "Opening Netflix"

    return "Cancelled"

def handle_games():
    exe = find_roblox()
    if exe:
        os.startfile(exe)
        return "Opening Roblox"
    return "No game found"

def handle_suggestions(cmd):
    """Handle game suggestion requests without launching games"""
    # Extract the actual request content
    _, content = separate_question_action(cmd)
    
    # Ask AI for game suggestions
    suggestion_prompt = f"""The user is asking for game suggestions with this request: "{content}"
    
    Please provide helpful game suggestions, recommendations, and ideas. 
    Do NOT launch or open any games - only provide suggestions and information.
    Focus on:
    - Popular games in different genres
    - Games suitable for different age groups
    - Free and paid game options
    - Games for different platforms (PC, mobile, console)
    - Educational and entertainment games
    - Latest trending games
    
    Be helpful and informative, but do NOT execute any game launching commands."""
    
    response = ask_ai(suggestion_prompt)
    return response

def handle_search(cmd):
    query = strip_leading_action(cmd, SEARCH_VERBS)
    
    # Use DDGS to get search results
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query or cmd, max_results=5))
        
        if not results:
            return f"No search results found for: {query or cmd}"
        
        # Extract search result information
        search_info = []
        for r in results:
            search_info.append(f"Title: {r.get('title', 'No title')}\nContent: {r.get('body', 'No content')}")
        
        # Ask AI to summarize the search results
        summary_prompt = f"""Please provide a comprehensive summary of the following search results for the query "{query or cmd}":

{chr(10).join(search_info)}

Focus on the most relevant and useful information. Do not include URLs or links. Provide a clear, concise summary that answers the user's query."""
        
        summary = ask_ai(summary_prompt)
        return f"Here's what I found about {query or cmd}:\n\n{summary}"
        
    except Exception as e:
        return f"Search error: {str(e)}"

def handle_voice_control(cmd):
    """Handle voice control commands"""
    cmd = cmd.lower().strip()
    
    if "mute" in cmd or "voice off" in cmd or "stop talking" in cmd:
        if is_voice_enabled():
            toggle_voice()
            return "Voice muted."
        return "Voice is already muted."
    
    elif "unmute" in cmd or "voice on" in cmd or "start talking" in cmd:
        if not is_voice_enabled():
            toggle_voice()
            return "Voice unmuted."
        return "Voice is already enabled."
    
    elif "toggle" in cmd:
        status = "enabled" if toggle_voice() else "disabled"
        return f"Voice {status}."
    
    elif "louder" in cmd:
        # Increase volume
        from voice_response import voice_responder
        voice_responder.set_voice_volume(min(1.0, voice_responder.engine.getProperty('volume') + 0.1))
        return "Volume increased."
    
    elif "quieter" in cmd:
        # Decrease volume
        from voice_response import voice_responder
        voice_responder.set_voice_volume(max(0.0, voice_responder.engine.getProperty('volume') - 0.1))
        return "Volume decreased."
    
    elif "faster" in cmd:
        # Increase speech rate
        from voice_response import voice_responder
        voice_responder.set_voice_rate(min(400, voice_responder.engine.getProperty('rate') + 20))
        return "Speech rate increased."
    
    elif "slower" in cmd:
        # Decrease speech rate
        from voice_response import voice_responder
        voice_responder.set_voice_rate(max(50, voice_responder.engine.getProperty('rate') - 20))
        return "Speech rate decreased."
    
    return "Voice commands: mute voice, unmute voice, toggle voice, speak louder/quieter, speak faster/slower"

def handle_tone_control(cmd):
    """Handle tone control commands including voice/text format toggle"""
    cmd = cmd.lower().strip()
    tone_mgr = get_tone_manager()
    
    # Handle voice/text format toggle
    if any(word in cmd for word in ["voice", "text", "format", "switch"]):
        if "voice" in cmd or "speak" in cmd:
            return "Switched to voice mode. I will now speak my responses aloud."
        elif "text" in cmd:
            return "Switched to text mode. I will now display responses as text only."
    
    if "list" in cmd or "available" in cmd:
        tones = tone_mgr.get_tone_list()
        tone_list = "\n".join([f"  {name}: {desc}" for name, desc in tones.items()])
        return f"Available tones:\n{tone_list}\n\nCurrent tone: {get_current_tone_name()}\n\nCurrent mode: {'Voice' if is_voice_enabled() else 'Text'}"
    
    # Extract tone name from command
    tone_patterns = [
        r"tone to (\w+)",
        r"set tone (\w+)",
        r"change tone to (\w+)",
        r"be more (\w+)",
        r"be less (\w+)",
        r"sound more (\w+)",
        r"sound less (\w+)",
        r"speak like (\w+)",
        r"talk like (\w+)",
        r"tone (\w+)",
        r"list tones",
        r"available tones",
    ]
    
    for pattern in tone_patterns:
        match = re.search(pattern, cmd)
        if match:
            tone_name = match.group(1)
            
            # Handle "more/less" modifiers
            if "more" in pattern or "sound" in pattern:
                # Map descriptive words to tone names
                tone_mapping = {
                    "professional": "professional",
                    "formal": "professional",
                    "business": "professional",
                    "friendly": "friendly",
                    "warm": "friendly",
                    "nice": "friendly",
                    "casual": "casual",
                    "relaxed": "casual",
                    "laid back": "casual",
                    "enthusiastic": "enthusiastic",
                    "excited": "enthusiastic",
                    "energetic": "enthusiastic",
                    "technical": "technical",
                    "precise": "technical",
                    "detailed": "technical",
                    "minimal": "minimal",
                    "brief": "minimal",
                    "short": "minimal",
                    "creative": "creative",
                    "artistic": "creative",
                    "imaginative": "creative",
                    "sassy": "sassy",
                    "witty": "sassy",
                    "cheeky": "sassy",
                }
                
                tone_name = tone_mapping.get(tone_name, tone_name)
            
            if set_ai_tone(tone_name):
                # Voice settings will be updated automatically in speak_response function
                return format_with_tone(f"Tone changed to {tone_name}", "confirmation")
            else:
                available = list(tone_mgr.get_tone_list().keys())
                return f"Tone '{tone_name}' not found. Available tones: {', '.join(available)}"
    
    return "Voice/Text commands: 'voice' or 'text' to switch format, 'list tones' to see options, 'change tone to [name]' to set personality"

def search_and_open_file(file_query):
    """Search for and open a specific file across all accessible locations"""
    import glob
    from pathlib import Path
    
    normalized_query = normalize(file_query).lower().strip()
    
    # Common search paths
    search_paths = [
        os.path.expanduser("~/Desktop"),
        os.path.expanduser("~/Documents"),
        os.path.expanduser("~/Downloads"),
        os.path.expanduser("~/Pictures"),
        os.path.expanduser("~/Videos"),
        os.path.expanduser("~/Music"),
        os.path.expanduser("~"),
        "C:\\Users\\",
        "C:\\Program Files\\",
        "C:\\Program Files (x86)\\",
    ]
    
    # File extensions to search for
    file_extensions = [
        "*.txt", "*.py", "*.js", "*.html", "*.css", "*.json", "*.xml", "*.yaml", "*.yml",
        "*.md", "*.pdf", "*.doc", "*.docx", "*.xls", "*.xlsx", "*.ppt", "*.pptx",
        "*.jpg", "*.jpeg", "*.png", "*.gif", "*.bmp", "*.mp4", "*.avi", "*.mp3",
        "*.exe", "*.msi", "*.zip", "*.rar", "*.7z"
    ]
    
    found_files = []
    
    # Search through all paths
    for search_path in search_paths:
        if not os.path.exists(search_path):
            continue
            
        try:
            for ext in file_extensions:
                pattern = os.path.join(search_path, "**", ext)
                for file_path in glob.glob(pattern, recursive=True):
                    file_name = os.path.basename(file_path)
                    file_name_lower = file_name.lower()
                    
                    # Check if file name contains the query
                    if normalized_query in file_name_lower or fuzz.partial_ratio(normalized_query, file_name_lower) > 70:
                        found_files.append({
                            'path': file_path,
                            'name': file_name,
                            'score': fuzz.partial_ratio(normalized_query, file_name_lower)
                        })
        except (PermissionError, OSError):
            continue
    
    # Sort by score (best match first)
    found_files.sort(key=lambda x: x['score'], reverse=True)
    
    if found_files:
        # Open the best match
        best_match = found_files[0]
        try:
            os.startfile(best_match['path'])
            return f"Opening '{best_match['name']}' from {best_match['path']}"
        except Exception as e:
            return f"Found '{best_match['name']}' but couldn't open it: {str(e)}"
    
    return f"File '{file_query}' not found in common locations"

def open_app(name):
    normalized_name = normalize(name)
    cleaned_name = strip_leading_action(normalized_name, OPEN_VERBS)

    expanded_name = os.path.expandvars(name.strip())
    if os.path.exists(expanded_name):
        os.startfile(expanded_name)
        return "Opening file"

    # Check if this is a file search request
    if any(keyword in cleaned_name.lower() for keyword in ["file", "document", "find", "search"]):
        return search_and_open_file(cleaned_name)

    if fuzzy_contains(cleaned_name, ["roblox"], threshold=85):
        exe = find_roblox()
        if exe:
            os.startfile(exe)
            return "Opening roblox"
        return "Roblox not found"

    for app, target in BUILT_IN_APPS.items():
        if app in cleaned_name or fuzz.partial_ratio(cleaned_name, app) >= 85:
            os.startfile(os.path.expandvars(target))
            return f"Opening {app}"

    if cleaned_name:
        try:
            os.startfile(os.path.expandvars(cleaned_name))
            return f"Opening {cleaned_name}"
        except OSError:
            pass

    # Try file search as last resort
    return search_and_open_file(cleaned_name)

# =========================
# MAIN
# =========================
def process_command(cmd, interactive=True, use_vision=False):
    normalized_cmd = normalize(cmd)

    #Hard-stop self-modification queries
    if is_self_modification_query(normalized_cmd):
        response = modify_self_code(cmd)
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    # Hard-stop voice control queries
    if is_voice_control_query(normalized_cmd):
        response = handle_voice_control(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    # Hard-stop tone control queries
    if ENABLE_TONE_SYSTEM and is_tone_control_query(normalized_cmd):
        response = handle_tone_control(cmd)
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    # Hard-stop app automation queries
    if any(keyword in normalized_cmd for keyword in ["fusion", "netflix", "youtube", "open ", "launch ", "start "]):
        response = handle_app_command(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    # Hard-stop vision queries first
    if is_vision_query(normalized_cmd):
        response = handle_vision(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response)
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    # Hard-stop common media intents before any chat fallback
    if is_listen_query(normalized_cmd) and not is_question(normalized_cmd):
        response = handle_listen()
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if is_game_query(normalized_cmd) and not is_question(normalized_cmd):
        response = handle_games()
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if is_watch_query(normalized_cmd) and not is_question(normalized_cmd):
        response = handle_watch(cmd, interactive=interactive)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if is_cad_model_query(normalized_cmd):
        response = format_cad_reply_cli(generate_cad_design(cmd))
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    # Detect intent using enhanced system
    intent, content, is_question = detect_enhanced_intent(cmd)
    
    # Handle suggestion intent - provide suggestions without launching games
    if intent == "suggest":
        response = handle_suggestions(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "helpful")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response
    
    if intent == "vision":
        response = handle_vision(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response)
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response
    
    if intent == "chat":
        response = ask_ai_maybe_with_camera(cmd, use_vision)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response)
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response
    
    if intent == "game":
        response = handle_games()
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "music":
        response = handle_music()
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "listen":
        response = handle_listen()
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "watch":
        response = handle_watch(cmd, interactive=interactive)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "cad":
        response = format_cad_reply_cli(generate_cad_design(cmd))
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "open":
        response = open_app(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "search":
        response = handle_search(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response, "confirmation")
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "news":
        response = handle_news(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response)
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "time":
        response = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response)
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    if intent == "self_modify":
        # Check if this is an autonomous improvement request
        if any(keyword in cmd.lower() for keyword in ["analyze yourself", "self improve", "make yourself better", "autonomous improvement", "improve autonomously"]):
            # Perform self-analysis first
            analysis = analyze_self_for_improvements()
            if ENABLE_TONE_SYSTEM:
                analysis = format_with_tone(analysis, "helpful")
            if ENABLE_VOICE_RESPONSE:
                speak_response(analysis)
            return analysis
        
        # Check if user wants to see improvement status
        if any(keyword in cmd.lower() for keyword in ["improvement status", "continuous improvement", "improvement queue"]):
            status = get_improvement_status()
            if ENABLE_TONE_SYSTEM:
                status = format_with_tone(status, "informative")
            if ENABLE_VOICE_RESPONSE:
                speak_response(status)
            return status
        
        # Check if user wants to implement a specific improvement
        if "implement" in cmd.lower() or "add" in cmd.lower():
            # Extract improvement description
            improvement_desc = cmd.lower()
            for prefix in ["implement", "add", "create", "build"]:
                if prefix in improvement_desc:
                    improvement_desc = improvement_desc.replace(prefix, "").strip()
                    break
            
            response = implement_autonomous_improvement(improvement_desc)
            if ENABLE_TONE_SYSTEM:
                response = format_with_tone(response, "confirmation")
            if ENABLE_VOICE_RESPONSE:
                speak_response(response)
            return response
        
        # Check if user wants to queue an improvement
        if "queue" in cmd.lower() or "improve later" in cmd.lower():
            improvement_desc = cmd.lower()
            for prefix in ["queue", "improve later"]:
                if prefix in improvement_desc:
                    improvement_desc = improvement_desc.replace(prefix, "").strip()
                    break
            
            response = queue_improvement(improvement_desc)
            if ENABLE_TONE_SYSTEM:
                response = format_with_tone(response, "confirmation")
            if ENABLE_VOICE_RESPONSE:
                speak_response(response)
            return response
        
        # Check if user wants to upgrade AI model
        if any(keyword in cmd.lower() for keyword in ["upgrade model", "change model", "recommend model", "ai model", "model upgrade"]):
            # Extract model name if specified
            new_model = None
            for model_name in ["llama3.1:8b", "qwen2.5:7b", "qwen2.5:1.5b", "mistral:7b", "phi3:mini"]:
                if model_name in cmd.lower():
                    new_model = model_name
                    break
            
            response = upgrade_ai_model(new_model)
            if ENABLE_TONE_SYSTEM:
                response = format_with_tone(response, "confirmation")
            if ENABLE_VOICE_RESPONSE:
                speak_response(response)
            return response
        
        # Default to manual self-modification
        response = modify_self_code(cmd)
        if ENABLE_TONE_SYSTEM:
            response = format_with_tone(response)
        if ENABLE_VOICE_RESPONSE:
            speak_response(response)
        return response

    response = ask_ai_maybe_with_camera(cmd, use_vision)
    if ENABLE_TONE_SYSTEM:
        response = format_with_tone(response)
    if ENABLE_VOICE_RESPONSE:
        speak_response(response)
    return response

def handle(cmd):
    print(process_command(cmd, interactive=True))

# =========================
# LOOP
# =========================
def main():
    print(f"{ASSISTANT_NAME} online...\nType 'exit' to quit\n")
    print(f"Text AI Model: {AI_MODEL}")
    
    if ENABLE_CAMERA:
        print(f"Vision Model: {VISION_MODEL}")
        print("Vision capabilities: ENABLED")
        print(f"Using camera index: {CAMERA_INDEX}\n")
    else:
        print("Vision capabilities: DISABLED\n")
    
    if ENABLE_CONTINUOUS_IMPROVEMENT:
        print("Continuous self-improvement: ENABLED")
        print(f"Improvement interval: {IMPROVEMENT_INTERVAL} seconds\n")
        # Start continuous improvement system
        start_continuous_improvement()
    else:
        print("Continuous self-improvement: DISABLED\n")

    while True:
        cmd = input(">>> ")
        if cmd.lower() in ["exit", "quit"]:
            break
        handle(cmd)

if __name__ == "__main__":
    main()
